from flask import Flask, request, render_template
from docx import Document
import pytesseract
import cv2
from deep_translator import GoogleTranslator
import os
import fitz
from pptx import Presentation
import webbrowser

pytesseract.pytesseract.tesseract_cmd = 'Tesseract-ocr\\tesseract.exe'

def Run_Model_CV(filepath,og_lang,dir_lang,model_lang,configg='--oem 3 --psm 6'):
  image = cv2.imread(filepath)
  image = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
  text = pytesseract.image_to_string(image, config=configg,lang=model_lang)
  chunks = [text[x:x+250] for x in range(0, len(text), 250)]
  translated_chunks = []
  i=1
  for chunk in chunks:
      translated_chunk = GoogleTranslator(source=og_lang, target=dir_lang).translate(chunk)
      print(str(i)+"/"+str(len(chunks)))
      i=i+1
      translated_chunks.append(translated_chunk)
  translated_text = ' '.join(translated_chunks)

  return translated_text

def extract_text_from_ppt(presentation_path):
    presentation = Presentation(presentation_path)
    all_text = []
    text_per_slide=[]
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text.strip()
                        if text:
                            all_text.append(text)
        text_per_slide.append(" ".join(all_text))
        all_text=[]

    return text_per_slide

# Usage example
def run_model_ppt(file_path,og_lang,dir_lang):
    presentation_path = file_path
    text = extract_text_from_ppt(presentation_path)
    text = '\n'.join(text)
    chunks = [text[x:x+250] for x in range(0, len(text), 250)]
    translated_chunks = []
    i=1
    for chunk in chunks:
        translated_chunk = GoogleTranslator(source=og_lang, target=dir_lang).translate(chunk)
        print(str(i)+"/"+str(len(chunks)))
        i=i+1
        translated_chunks.append(translated_chunk)
    translated_text = ' '.join(translated_chunks)

    return translated_text

def run_translate_pdf(filepath,og_lang,dir_lang):
  pdf_document = fitz.open(filepath)
  num_pages = pdf_document.page_count
  text=''
  for page_num in range(num_pages):
    page = pdf_document.load_page(page_num)
    text += page.get_text()
  chunks = [text[x:x+250] for x in range(0, len(text), 250)]
  translated_chunks = []
  i=1
  for chunk in chunks:
      translated_chunk = GoogleTranslator(source=og_lang, target=dir_lang).translate(chunk)
      print(str(i)+"/"+str(len(chunks)))
      i=i+1
      translated_chunks.append(translated_chunk)
  translated_text = ' '.join(translated_chunks)

  return translated_text

def run_model_txt(file_path,og_lang,dir_lang):
    file = open(file_path, 'r' , encoding='utf-8')
    text= file.read()
    chunks = [text[x:x+250] for x in range(0, len(text), 250)]
    translated_chunks = []
    i=1
    for chunk in chunks:
        translated_chunk = GoogleTranslator(source=og_lang, target=dir_lang).translate(chunk)
        print(str(i)+"/"+str(len(chunks)))
        i=i+1
        translated_chunks.append(translated_chunk)
    translated_text = ' '.join(translated_chunks)

    return translated_text



app = Flask(__name__)



@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')


@app.route('/display', methods=['POST'])
def display():
    file_path = request.form['file_path']
    og_lang = request.form['og_lang']
    dir_lang = request.form['dir_lang']
    
    model_lang = request.form['model_lang']
    _, file_extension = os.path.splitext(file_path)

    if file_extension.lower() in ['.jpg', '.png', '.gif', '.bmp']:
        return render_template('display.html', data=Run_Model_CV(file_path,og_lang,dir_lang,model_lang), file_type="PythonString")     

    elif file_extension.lower() in ['.pdf']:
        return render_template('display.html', data=run_translate_pdf(file_path,og_lang,dir_lang), file_type="PythonString")   
    
    elif file_extension.lower() in ['.pptx']:
        return render_template('display.html',data=run_model_ppt(file_path,og_lang,dir_lang), file_type = 'PythonString')
    
    elif file_extension.lower() in ['.txt','.text']:
        return render_template('display.html',data=run_model_txt(file_path,og_lang,dir_lang), file_type = 'PythonString')
    
    elif file_extension.lower() in ['.docx']:
        doc = Document(file_path)
        data = []
        Paras=doc.paragraphs;
        par=1
        for paragraph in Paras:
            print(str(par)+"/"+str(len(Paras)))
            text = paragraph.text
            chunks = [text[x:x+250] for x in range(0, len(text), 250)]
            translated_chunks = []
            for chunk in chunks:
                translated_chunk = GoogleTranslator(source=og_lang, target=dir_lang).translate(chunk)
                translated_chunks.append(translated_chunk)
            translated_text = ' '.join(translated_chunks)
            text = translated_text
            font_size = paragraph.style.font.size.pt if paragraph.style.font.size else None
            indentation = paragraph.style.paragraph_format.first_line_indent.pt if paragraph.style.paragraph_format.first_line_indent else None
            par=par+1
            data.append({'text': text, 'font_size': font_size, 'indentation': indentation})
    
        return render_template('display.html', data=data, file_type='')
    
    else:
        return render_template('display.html',data='Unsupported File Formate', file_type='PythonString')

        
if __name__ == "__main__":
    webbrowser.open('http://localhost:5000')
    app.run(port=5000,debug=True)

# pyinstaller --onefile --add-data "Test_Flask/templates:templates" --add-data "Test_Flask/static:static" --add-data "Tesseract-ocr:Tesseract-ocr" Test_Flask/Test_Flask.py